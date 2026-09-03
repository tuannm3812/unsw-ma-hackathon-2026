#!/usr/bin/env python3
"""Regenerate docs/presentation/slides_draft.pptx from the deck brief content.

Single source of truth for the editable draft deck: slide content, big-number
callouts, chart placeholders (named after the notebook section to screenshot)
and the timed speaker scripts (as presenter notes). Design mirrors the deck
brief page: warm paper ground, deep teal + amber accent, Georgia display type.
Sync contract: docs/presentation/deck_content.md is the source of truth for
wording; this script mirrors it (update both together). Chart assets live in
docs/presentation/charts/ and are the source of truth for exhibits. The build
FAILS if any required chart is missing; pass --scaffold to emit labelled
placeholders instead. Requires python-pptx (optional tooling, not a pipeline
dependency): pip install python-pptx
Re-run after any wording change: python3 scripts/build_slides_draft.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import os

# Slide -> chart image (extracted from the Kaggle-executed notebooks; the
# few-cluster table is typeset from the verified v15 SS7.2 output). Override
# with CHART_DIR env var; missing files fall back to the labelled placeholder.
CHART_DIR = os.environ.get("CHART_DIR", "docs/presentation/charts")
CHARTS = {
    3: "data_split.png",        # chronological-split schematic (verified counts)
    4: "period_24h.png",        # 24h-funded share by period (rebuilt, exact shares)
    5: "sector.png",            # average funding speed by sector
    6: "few_cluster_table.png", # SS7.2 within-region few-cluster screen (typeset)
    7: "topics.png",            # topic mean speeds, semantic labels (8 NMF topics)
    12: "region.png",           # appendix: region speeds (house style, verified values)
    13: "shap_top15.png",       # appendix: SHAP top-15, human feature names
    14: "correlations.png",     # appendix: the exact 10x correlation basis
}

# Deck palette: navy ink + viridis blue + team yellow #FFDD04 (yellow is a
# shape/chip colour only - never text on a light ground)
PAPER = RGBColor(0xFB, 0xFA, 0xF6); INK = RGBColor(0x1C, 0x23, 0x33)
TEAL = RGBColor(0x1C, 0x23, 0x33); CREAM = RGBColor(0xFB, 0xFA, 0xF6)
AMBER = RGBColor(0xFF, 0xDD, 0x04); MUTED = RGBColor(0x6E, 0x72, 0x78)
PANEL = RGBColor(0xEF, 0xED, 0xE6); LINE = RGBColor(0xD9, 0xD6, 0xCC)
BLUE = RGBColor(0x31, 0x68, 0x8E)
# DM Sans matches the chart typeface (install docs/presentation/fonts/*.ttf
# on any machine that opens this deck, or PowerPoint silently falls back).
DISPLAY = "Georgia"; BODY = "DM Sans"; MONO = "Courier New"

# Academic figure captions shown under each embedded exhibit.
FIGS = {
    3:  (1, "Chronological train/test split.",
         "The 278,887 test loans stayed sealed during training - the models never see the future they are scored on. Counts: analysis_summary.json (boundary 2024-01-01)."),
    4:  (2, "Share of loans funded within 24 hours, by analysis period.",
         "The drop is 16 points; the 'recovery' is -0.3. Exact shares 46.0% / 30.3% / 30.0% (589,823 / 298,549 / 565,474 loans). EDA notebook \u00a74."),
    5:  (3, "Mean funding speed by sector.",
         "Sanitation loans fund ~13x faster than Clothing (0.9 vs 12.1 days). Computed from the raw data replicating EDA \u00a75 grouping; every count matches the executed notebook."),
    6:  (4, "Average within-region family-framing slope (duration model).",
         "Against the few-cluster screen - a conservative heuristic (t(1) critical value 12.7), able to downgrade a claim but never certify one - nothing clears it. Typeset from executed \u00a77.2 (v15); 4-dp values in analysis_summary.json."),
    7:  (5, "Mean funding speed by story theme.",
         "Water and sanitation stories fund in under 2 days; group farming waits ~2 weeks - themes track what the loan is FOR. EDA \u00a78 printed means (TF-IDF + NMF, 8 topics, 20K sample)."),
    8:  None,
    12: (6, "Mean funding speed by region.",
         "North America is Haiti alone (n = 7,559) - fastest region, and exactly why single-cluster inference is off the table. Computed from the raw data (counts verified vs executed notebook)."),
    13: (7, "Top 15 features by mean |SHAP| value, boosted model.",
         "Ten structural features outrank the first narrative one. Rebuilt from modeling \u00a78 printed values (2,000-loan holdout sample)."),
    14: (8, "Correlation with funding speed: structural vs narrative features.",
         "Loan amount correlates ~7x stronger with speed than the best narrative signal (0.429 vs 0.058). EDA \u00a79 printed correlation table."),
}

DEFAULT_SOURCE = ("reports/generated_full_dataset/association_summary.txt (authoritative numbers); docs/presentation/deck_content.md SS2 numbers table")

CLOSING_LINES = ["\u201cThe story barely registers.", "The structure carries the signal.\u201d"]

SOURCES = {
    3: "Split schematic drawn from analysis_summary.json split sizes (train 1,174,953 / holdout 278,887; boundary 2024-01-01).",
    4: "Chart rebuilt from EDA SS4 printed shares (0.460/0.303/0.300) and period counts; association_summary.txt audit trail.",
    5: "Sector means computed from data/Kiva_Loans.pkl replicating EDA SS5 grouping (counts verified vs executed notebook). Gender medians: EDA SS4 printout.",
    6: "Table typeset from the executed modeling notebook SS7.2 printout (v15); authoritative few-cluster values: analysis_summary.json within_region_slopes.",
    7: "Chart rebuilt from EDA SS8 printed topic means and top-words (8 NMF topics).",
    8: "Classifier metrics: analysis_summary.json binary_classifier. Few-cluster screen: SS7.2 / within_region_slopes.",
    12: "Region means computed from data/Kiva_Loans.pkl replicating EDA SS5 grouping (counts verified). Pooled definitions + p-values: association_summary.txt.",
    13: "Chart rebuilt from modeling SS8 printed SHAP top-15 values (v14/v15 runs).",
    14: "Chart rebuilt from EDA SS9 printed correlation table; decile curve: EDA SS9 figure.",
}

SLIDES = [
 ("Beyond a Good Story",
  "When — and for whom — does persuasive loan language actually speed up funding?",
  None, ["Team Cultural Blend  ·  Kiva loan data  ·  1.45 million real loans (2016–2025)"], None,
  "(Sophia, ~30s) Good morning — we're Cultural Blend: I'm Sophia, and my teammate Tuan will take you through what survives when we test it hard. Kiva is built on stories — every loan page leads with one, the way a landing page leads with copy. We asked 1.45 million real loans a single question: does the story actually move the money?",
  "UNSW MARKETING ANALYTICS HACKATHON · FINAL"),
 ("The question", None, None,
  ["When a loan's story leans on family, competence, or urgency — does it fund faster?",
   "Does the answer depend on WHO is asking and WHEN?",
   "Why it matters: framing is the one thing a platform can coach — loan size, sector and geography can't be rewritten after the fact."],
  None,
  "(Sophia, ~50s) Why should a marketing audience care? Because on Kiva, the lender is the customer and the loan page is the product page. Loan size, sector, geography — fixed at listing. The story is the one element a platform can coach, test, and optimise — classic conversion territory. So: when a story leans on family, competence, or urgency, does the loan fund faster? And does that depend on who's asking, and when? That's a testable claim — so we tested it, hard.",
  "MOTIVATION"),
 ("How we stress-tested our findings", None, None,
  ["Predictive claims: trained only on the past, tested only on loans posted 2024–2025 — no peeking at the future.",
   "Framing claims: every 'significant' result re-tested with country-clustered standard errors, then a harsher few-cluster reference where a result rests on a handful of countries.",
   "Most headline-looking results did NOT survive — that is the point of the method.",
   "SHAP importance shown as complementary predictive evidence only — it cannot confirm or refute a statistical finding."],
  "chronological train/test split · counts from analysis_summary.json",
  "(Sophia, ~65s) Two disciplines before any findings — because in a dataset this size it is dangerously easy to find things that aren't there. For prediction: train only on the past, score only on loans posted in 2024-25 — the models never see the future they're graded on. For the framing claims: every 'significant' result had to survive country-clustered standard errors — ten thousand loans from one country are not ten thousand independent customers — and where a result rested on just a couple of countries, an even harsher few-cluster screen on top. Most headline-looking results did not survive. Hold that thought, because it decides everything that follows.",
  "METHOD"),
 ("A marketplace that hasn't recovered", None,
  "46% → 30% → 30%|funded within 24 hours",
  ["Pre-pandemic, almost half of all loans funded within a day.",
   "Since 2020: under a third — and through the end of the data (2025) it has NOT recovered. Persistence to date, not proof it never will.",
   "589,823 loans before vs. 565,474 after — not a small, noisy sample."],
  "EDA notebook §4 categorical trends (period chart)",
  "(Sophia, ~65s) Start with the market, because every other number lives inside it. Before the pandemic, 46% of loans funded within 24 hours — nearly half converting same-day. Since 2020 it's been under a third — and through 2025, the 'recovery' is minus 0.3 points. With over half a million loans on each side of that divide, this isn't noise: it's a structurally slower, more selective marketplace. For a marketer that reframes the whole job — when customers get pickier, knowing what actually drives conversion matters more, not less. Four years without a recovery is also why our recommendations start with acting rather than waiting for the market to fix itself.",
  "FINDING 1 · THE MARKET"),
 ("Structure beats story", None,
  "2.3 vs 7.7 days|median funding speed — female- vs male-posted loans",
  ["Loan amount + repayment terms: linked to funding speed ~10x more strongly than any single narrative choice.",
   "Sector alone spans well over an order of magnitude in speed (0.9–12.1 days); country gaps are wider still — 0.2 to 20.5.",
   "None of this is causal — but it dwarfs anything the words do."],
  "EDA notebook §5 categorical features · §9 feature correlations",
  "(Sophia, ~65s) And what drives it is structure. Loan amount and repayment terms correlate with funding speed roughly ten times more strongly than any narrative signal. Sector alone spans an order of magnitude — sanitation stories fund in under a day, clothing takes twelve. And the starkest gap in the data: loans posted by women fund in a median of 2.3 days, by men 7.7 — three times longer. None of this is causal, but it's enormous, it's stable, and it dwarfs anything the words do. Which brings us to the question we actually came here to answer. Tuan —",
  "FINDING 2 · STRUCTURE"),
 ("What survives scrutiny", None,
  "No narrative result|is robust enough to support a recommendation",
  ["Urgency: looked universal (p<0.001) — collapses under country clustering (p≈0.44). Not recommended.",
   "Family: corrected test → Middle East (Palestine+Yemen) & Central America (Honduras+Nicaragua) same direction in every fit — but 2 countries each; few-cluster t(1) screen bar is 12.7, not 1.96: p≈0.06–0.21. A hypothesis, not a finding.",
   "Sentiment: positive tone <-> slower funding, but significance flips between specifications — genuinely open.",
   "SHAP: no family / agency / urgency feature in the model's top 15 — complementary evidence, not confirmation."],
  "modeling notebook §7.2 within-region slopes · §7.1 check",
  "(Tuan, ~1m50s) Thanks, Sophia. So: does the story matter? Our honest answer: no narrative result is robust enough across specifications to support a recommendation. Urgency language looked like a universal win — significant at p below 0.001. Cluster by country, and it collapses to 0.44. Gone. Family framing — and here our own first version got it wrong: we tested whether regions differ from Africa, which is not the same as whether family framing helps within a region. Corrected, two pooled categories — Palestine plus Yemen, and Honduras plus Nicaragua — do show faster funding in every fit we ran. But each rests on exactly two countries, and against a deliberately harsh few-cluster screen — a conservative heuristic, not calibrated inference: a t distribution with one degree of freedom, critical value 12.7, not 1.96 — neither is significant: p between 0.06 and 0.21. So we report a hypothesis worth testing, not a finding. Sentiment is the honest illustration: it survives country clustering in one of our two specifications and not the other — genuinely open, robust in neither direction. We'd rather report no robust evidence than one exciting result we can't defend — because a recommendation you'd ship to real borrowers deserves that bar.",
  "FINDING 3 · THE HEADLINE"),
 ("Beyond keywords", None,
  "1.5 → 13.5 days|mean funding speed across story topics — a >9x swing",
  ["Topic modeling on the descriptions (TF-IDF + NMF, 8 topics — not keyword counts) surfaces coherent themes: sanitation, clean water, pig raising, family business, smallholder farming.",
   "The largest single gap anywhere in the analysis — but a topic mostly encodes what the loan is FOR. Structure again, not persuasion."],
  "EDA notebook §8 topic modeling",
  "(Tuan, ~50s) One more layer before the recommendations — we went past keyword counting. Topic modelling finds eight coherent story themes, and mean funding speed swings ninefold across them — descriptively, not as a tested effect: sanitation and clean-water stories in under two days, group farming closer to two weeks. But notice what a theme mostly encodes: what the loan is FOR. Structure again — not persuasion — and it's why the structural review we're about to recommend should cover what loans are FOR, not just which sector they sit in. So what should Kiva actually do with all of this?",
  "FINDING 4 · TOPICS"),
 ("What this means in practice", None, None,
  ["DON'T ship writing tips — a platform-wide 'add urgency' nudge would be built on a result that doesn't survive testing.",
   "DO run a country-stratified A/B test of family framing in exactly those 4 markets — that's how a hypothesis becomes a decision.",
   "DO review the structural gaps (programme / sector / region / gender) — they are ~10x the size of any wording effect.",
   "PROTOTYPE, then re-validate: the 24h classifier ranks well among eventually-funded loans (AUC ≈ 0.90, strictly future data) — but expired/withdrawn listings aren't in the data, so an early-warning flag first needs all posted listings, a defined target, and retraining on that population.",
   "FOR FIELD PARTNERS (and borrowers): don't optimise wording — no robust evidence it pays. Timing is a partner capital-cycle matter, not a borrower wait: 96% of loans are already disbursed before the page goes live (median 24 days), and a bigger ask sits on the platform longer (~2 vs ~19 days). Still NOT a recommendation to ask for less."],
  None,
  "(Tuan, ~1m25s) Three moves for the platform — and notice the verbs, because the evidence sets them. One: don't ship writing tips — a firm don't, because a robust null is the one thing we did find; for the content team that's a build spared, and for borrowers it's not being coached into copy we found no robust evidence for. Two — an experiment, not an action, because there we have a pattern we could not confirm: for the growth team and the field partners in exactly four markets — Palestine, Yemen, Honduras, Nicaragua — run the country-stratified A/B test: family-framing prompt versus standard at listing. Test before you ship. Three — review, not change, because these are associations and not causes: for the product team, the structural gaps are the real levers — how the consistently slower programmes, sectors and regions get surfaced and supported — because those gaps are ten times any wording effect. And the classifier stays a prototype, not a deployment — a retrospective ranking prototype among funded loans — AUC 0.90 on strictly future data — so data science retrains on all listings, including expired ones, first. [advance to the borrower page] (Tuan, ~40s) And the borrower page — expectations, not tactics. The timing story belongs to field partners: a bigger ask sits on the platform longer — two days for the smallest loans against nineteen for the largest — and since 96% of these loans are already disbursed before the page even goes live, that is a partner capital-cycle fact, not a borrower's waiting time. Which is why we are not telling anyone to ask for less — this page sets expectations by profile, amount and sector, and none of them is a lever to game.",
  "RECOMMENDATIONS"),
 ("What this can't tell us", None, None,
  ["Association, never causation — no borrower was randomly assigned a writing style.",
   "Measures how fast a FUNDED loan funds — not whether a loan gets funded at all (expired/withdrawn listings never enter the data).",
   "Framing measured with transparent, simple rules — not every nuance of persuasive writing."],
  None,
  "(Tuan, ~35s) Three honest limits. Association, never causation — nobody randomly assigned writing styles. We measure how fast funded loans fund — not whether a loan funds at all. And our framing measures are transparent, simple rules — not every nuance of persuasion. We'd rather you know exactly what this can and cannot say — that's what makes the parts we do claim worth trusting.",
  "LIMITS"),
 ("", None, None, [], None,
  "(Both, ~20s) (Tuan) In this data, the story barely registers — the structure carries the signal. (Sophia) And testing hard enough to say, honestly, that there is no robust evidence for the story — that's worth more to a platform than a good-sounding tip. Thank you — we're happy to take your questions.",
  "CLOSING"),
 # ---- Appendix: Q&A backup, never presented in the 10 minutes ----
 ("Appendix", "Q&A backup — not part of the 10-minute presentation.", None, [], None,
  "Divider only. Everything after this slide exists to be pulled up during Q&A if a question calls for it.",
  "APPENDIX"),
 ("The two-country problem, in full", None, None,
  ["Pooled categories: Middle East = Palestine + Yemen (14,946 loans); Central America = Honduras + Nicaragua (59,391); North America = Haiti alone.",
   "Few-cluster screen p (duration / 24h / notebook fit): ME 0.12 / 0.21 / 0.08 · CA 0.06 / 0.14 / 0.07 — same direction in every fit, none significant; t(1) critical value 12.7.",
   "North America: single country — between-country uncertainty not estimable at all.",
   "The screen is a conservative heuristic: it can downgrade a claim, never certify one. Next step: country-stratified A/B test in those four markets."],
  "EDA notebook §5 · regional speed context for the pooled categories",
  "Backup for Slide 6 probes: exact pooled definitions, all p-values, the single-country boundary, and the honest status of the screen.",
  "APPENDIX · A1"),
 ("Predictive weight ≠ statistical robustness", None, None,
  ["SHAP importance from the boosted forecasting model: structure fills the top ranks (amount, term, period, size band, sector, gender).",
   "No family, agency, or urgency feature reaches the top 15; sentiment is 11th despite its disputed significance.",
   "Complementary predictive evidence only — a different model, no region interactions; it cannot confirm or refute any coefficient's sign or uncertainty."],
  "modeling notebook §8 · SHAP top-15",
  "Backup for SHAP questions: what the ranking shows, and exactly what it cannot corroborate.",
  "APPENDIX · A2"),
 ("The '~10x' comparison, exactly", None, None,
  ["Correlation with funding speed: loan amount (log) r = 0.43, repayment term r = 0.28 — vs competence 0.058, family 0.019, urgency 0.010.",
   "Loan-amount deciles also rise almost monotonically (~2 days for the smallest asks, ~18 for the largest) — monotonicity, not the 10x basis itself.",
   "Gender: female-posted median 2.3 days vs male 7.7 (male coefficient +0.430 survives HC3 AND country clustering, p < 0.0001 — still associational)."],
  "EDA notebook §9 · correlation table (the 10x basis)",
  "Backup for Slide 5 probes: the amount gradient and the gender gap's robustness status.",
  "APPENDIX · A3"),
 ("24h classifier — operating detail", None, None,
  ["Holdout (posted 2024-01-01 onward): 278,887 loans; 87,466 funded within 24h vs 191,421 not.",
   "ROC AUC 0.8997 · average precision 0.8301 · Brier 0.1156 · accuracy 0.840 — structure does the work: no framing feature in the SHAP top 10 (sentiment is 11th).",
   "Boundary: negative class = 'eventually funded, but not within 24h'. Expired/withdrawn listings never enter the data — a retrospective ranking prototype among eventual funders.",
   "Path to deployment: all posted listings incl. expired/withdrawn outcomes -> define operational target + censoring window -> retrain, validate -> threshold, calibration, fairness, prospective test."],
  None,
  "Backup for classifier questions: exact metrics, class balance, and the population boundary stated before anyone asks.",
  "APPENDIX · A4"),
]


def build(out_path: str) -> None:
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def bg(slide, color):
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color

    def box(slide, x, y, w, h):
        b = slide.shapes.add_textbox(x, y, w, h); b.text_frame.word_wrap = True
        return b.text_frame

    def style(p, size, color, font=BODY, bold=False, italic=False, align=None):
        p.font.size = Pt(size); p.font.color.rgb = color; p.font.name = font
        p.font.bold = bold; p.font.italic = italic
        if align: p.alignment = align

    def rect(slide, x, y, w, h, fill, line=None, rounded=False):
        sh = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h)
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
        if line: sh.line.color.rgb = line; sh.line.width = Pt(0.75)
        else: sh.line.fill.background()
        sh.shadow.inherit = False
        return sh

    for i, (title, sub_, big, bullets, chart, script, eyebrow) in enumerate(SLIDES, 1):
        sl = prs.slides.add_slide(blank)
        bg(sl, TEAL if i in (1, 10, 11) else PAPER)
        if i == 11:
            rect(sl, Inches(0.9), Inches(3.0), Inches(1.6), Pt(4), AMBER)
            tf = box(sl, Inches(0.9), Inches(3.25), Inches(11.5), Inches(1.2))
            tf.text = title; style(tf.paragraphs[0], 54, CREAM, DISPLAY, bold=True)
            tf = box(sl, Inches(0.9), Inches(4.5), Inches(11.5), Inches(0.6))
            tf.text = sub_; style(tf.paragraphs[0], 18, RGBColor(0xB8, 0xCE, 0xD6), BODY, italic=True)
        elif i == 1:
            rect(sl, Inches(0.9), Inches(2.05), Inches(1.6), Pt(4), AMBER)
            tf = box(sl, Inches(0.9), Inches(1.35), Inches(11.5), Inches(0.5))
            tf.text = eyebrow; style(tf.paragraphs[0], 13, AMBER, MONO)
            tf = box(sl, Inches(0.85), Inches(2.3), Inches(11.8), Inches(2.2))
            tf.text = title; style(tf.paragraphs[0], 66, CREAM, DISPLAY, bold=True)
            tf = box(sl, Inches(0.9), Inches(4.35), Inches(9.8), Inches(1.2))
            tf.text = sub_; style(tf.paragraphs[0], 22, RGBColor(0xB8, 0xCE, 0xD6), BODY, italic=True)
            tf = box(sl, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.6))
            tf.text = bullets[0]; style(tf.paragraphs[0], 15, RGBColor(0x8F, 0xB0, 0xBB), MONO)
        elif i == 10:
            rect(sl, Inches(5.87), Inches(1.7), Inches(1.6), Pt(4), AMBER)
            tf = box(sl, Inches(1.2), Inches(2.5), Inches(10.9), Inches(2.6))
            for j, line in enumerate(CLOSING_LINES):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = line; style(p, 42, CREAM, DISPLAY, bold=True, align=PP_ALIGN.CENTER)
            tf = box(sl, Inches(1.2), Inches(5.3), Inches(10.9), Inches(0.7))
            tf.text = "Thank you — questions."
            style(tf.paragraphs[0], 20, RGBColor(0xB8, 0xCE, 0xD6), BODY, align=PP_ALIGN.CENTER)
            tf = box(sl, Inches(1.2), Inches(6.9), Inches(10.9), Inches(0.4))
            tf.text = "Team Cultural Blend · UNSW Marketing Analytics Hackathon 2026"
            style(tf.paragraphs[0], 12, RGBColor(0x8F, 0xB0, 0xBB), MONO, align=PP_ALIGN.CENTER)
        else:
            tf = box(sl, Inches(0.75), Inches(0.42), Inches(9.0), Inches(0.35))
            tf.text = eyebrow; style(tf.paragraphs[0], 11.5, BLUE, MONO)
            chip = rect(sl, Inches(12.15), Inches(0.4), Inches(0.72), Inches(0.38), TEAL, rounded=True)
            ctf = chip.text_frame; ctf.text = f"{i}/10" if i <= 10 else f"A{i - 11}"
            ctf.margin_top = Emu(0); ctf.margin_bottom = Emu(0)
            style(ctf.paragraphs[0], 12, AMBER, MONO, align=PP_ALIGN.CENTER)
            ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = box(sl, Inches(0.7), Inches(0.78), Inches(11.3), Inches(1.0))
            tf.text = title; style(tf.paragraphs[0], 36, INK, DISPLAY, bold=True)
            rect(sl, Inches(0.75), Inches(1.62), Inches(1.35), Pt(3.5), AMBER)
            y = 1.95
            chart_top = y   # exhibits sit beside callout+bullets, full column height
            if big:
                num, label = big.split("|")
                tf = box(sl, Inches(0.75), Inches(y), Inches(7.35 if chart else 12.1), Inches(1.0))
                p = tf.paragraphs[0]; p.text = num.strip() + "  "
                style(p, 40, INK, DISPLAY, bold=True)
                r = p.add_run(); r.text = label.strip()
                r.font.size = Pt(16); r.font.color.rgb = MUTED; r.font.name = BODY; r.font.bold = False
                y += 1.15
            img = None
            if chart and CHART_DIR and i in CHARTS:
                cand = Path(CHART_DIR) / CHARTS[i]
                if cand.is_file():
                    img = cand
            bw = 12.1 if not chart else 7.35
            tf = box(sl, Inches(0.75), Inches(y), Inches(bw), Inches(6.9 - y))
            for j, item in enumerate(bullets):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                r1 = p.add_run(); r1.text = "▪  "; r1.font.color.rgb = AMBER; r1.font.size = Pt(16); r1.font.name = BODY
                r2 = p.add_run(); r2.text = item; r2.font.size = Pt(16.5); r2.font.color.rgb = INK; r2.font.name = BODY
                p.space_after = Pt(12); p.line_spacing = 1.15
            if chart and img is not None:
                # exhibits are authored at their physical on-slide size (300
                # dpi) - place at NATIVE size so type stays uniform across
                # charts; downscale only if one overflows its column/height
                pic = sl.shapes.add_picture(str(img), Inches(8.2), Inches(chart_top))
                max_w = Inches(4.55); max_h = Inches(6.55 - chart_top)
                # FAIL-FAST native-size contract (Codex round 11): every chart
                # is authored at its placed size; an oversized exhibit is a
                # chart-authoring bug and must break the build, never be
                # silently rescaled into illegibility.
                if pic.width > max_w * 1.005 or pic.height > max_h * 1.005:
                    raise SystemExit(
                        f"slide {i}: exhibit {img.name} authored "
                        f"{pic.width / 914400:.2f}x{pic.height / 914400:.2f}in exceeds its "
                        f"{max_w / 914400:.2f}x{max_h / 914400:.2f}in slot - re-author the "
                        "figure at its placed size (scripts/build_charts.py)"
                    )
                pic.left = Inches(8.2) + max(0, (Inches(4.55) - pic.width) // 2)
                fig_no, fig_name, fig_notes = FIGS[i]  # chart slides only; None entries never reach here
                cap = box(sl, Inches(8.2), Inches(6.62), Inches(4.7), Inches(0.75))
                cp = cap.paragraphs[0]
                r1 = cp.add_run(); r1.text = f"Figure {fig_no}. "
                r1.font.size = Pt(10); r1.font.bold = True; r1.font.color.rgb = INK; r1.font.name = BODY
                r2 = cp.add_run(); r2.text = fig_name
                r2.font.size = Pt(10); r2.font.color.rgb = INK; r2.font.name = BODY
                cp2 = cap.add_paragraph()
                r3 = cp2.add_run(); r3.text = "Notes: " + fig_notes
                r3.font.size = Pt(9); r3.font.color.rgb = MUTED; r3.font.name = BODY
            elif chart:
                ph = rect(sl, Inches(9.2), Inches(y), Inches(3.5), Inches(6.55 - y), PANEL, line=LINE, rounded=True)
                ptf = ph.text_frame; ptf.word_wrap = True; ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
                ptf.text = "CHART GOES HERE"
                style(ptf.paragraphs[0], 13, TEAL, MONO, bold=True, align=PP_ALIGN.CENTER)
                p2 = ptf.add_paragraph(); p2.text = "\nscreenshot from:\n" + chart
                style(p2, 11, MUTED, BODY, align=PP_ALIGN.CENTER)
            tf = box(sl, Inches(0.75), Inches(7.05), Inches(9.0), Inches(0.35))
            tf.text = "Cultural Blend · Kiva 2016–2025 · 1.45M loans"
            style(tf.paragraphs[0], 10, MUTED, MONO)
        src_note = SOURCES.get(i, DEFAULT_SOURCE)
        sl.notes_slide.notes_text_frame.text = script + "\n\n[Sources] " + src_note
    prs.save(out_path)
    # Post-build self-check: speaker-tagged script openers and a [Sources]
    # block on every notes page (fails the build on any drift).
    check = Presentation(out_path)
    for idx, slide in enumerate(check.slides):
        i = idx + 1
        expected = SLIDES[idx][5] + "\n\n[Sources] " + SOURCES.get(i, DEFAULT_SOURCE)
        notes = slide.notes_slide.notes_text_frame.text
        if notes != expected:
            raise SystemExit(
                f"slide {i}: notes drifted from source structures\n"
                f"  expected: {expected[:60]!r}...\n  got:      {notes[:60]!r}...")
    print("wrote", out_path, Path(out_path).stat().st_size,
          "bytes; exact notes self-check passed for", len(SLIDES), "slides")


if __name__ == "__main__":
    import sys
    scaffold = "--scaffold" in sys.argv
    if not scaffold and CHART_DIR:
        missing = [f for f in CHARTS.values() if not (Path(CHART_DIR) / f).is_file()]
        if missing:
            raise SystemExit(f"missing chart assets in {CHART_DIR}: {missing} (pass --scaffold to build placeholders)")
    build("docs/presentation/slides_draft.pptx")
